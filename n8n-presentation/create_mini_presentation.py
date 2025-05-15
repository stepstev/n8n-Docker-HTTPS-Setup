import os
from datetime import datetime
import time
import random
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

# Configuration des chemins
OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_FILE = os.path.join(OUTPUT_DIR, "n8n_Global_Presentation.pptx")
TEMPLATE_FILE = r"D:\00-Conception_AI\__Automations_n8n\n8n-presentation\ThinkinG_modele.pptx"
GITHUB_LINK = "https://github.com/votre-username/n8n-docker-docs"

# Définition des couleurs pour le texte sur fond sombre
TITLE_COLOR = RGBColor(240, 240, 240)  # Blanc cassé pour les titres
TEXT_COLOR = RGBColor(220, 220, 220)   # Gris clair pour le texte principal
ACCENT_COLOR = RGBColor(86, 180, 233)  # Bleu clair pour l'accentuation

def generate_unique_filename(base_path):
    """Génère un nom de fichier unique en cas de conflit."""
    if not os.path.exists(base_path):
        return base_path
    
    # Séparer le nom et l'extension
    directory, filename = os.path.split(base_path)
    name, ext = os.path.splitext(filename)
    
    # Créer un nouveau nom avec timestamp et nombre aléatoire
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    random_num = random.randint(1000, 9999)
    new_filename = f"{name}_{timestamp}_{random_num}{ext}"
    
    return os.path.join(directory, new_filename)

def save_presentation_safely(prs, filepath):
    """Sauvegarde la présentation avec gestion des erreurs de permission."""
    max_attempts = 3
    wait_time = 2  # secondes
    
    for attempt in range(max_attempts):
        try:
            if os.path.exists(filepath):
                print(f"Le fichier {filepath} existe déjà.")
                try:
                    # Tenter de le supprimer
                    os.remove(filepath)
                    print("Ancien fichier supprimé avec succès.")
                except PermissionError:
                    # Si on ne peut pas le supprimer, générer un nom unique
                    filepath = generate_unique_filename(filepath)
                    print(f"Impossible de supprimer l'ancien fichier. Utilisation d'un nouveau nom: {filepath}")
            
            # Essayer de sauvegarder la présentation
            prs.save(filepath)
            print(f"Présentation sauvegardée avec succès: {filepath}")
            return filepath
            
        except PermissionError as e:
            print(f"Erreur de permission (tentative {attempt+1}/{max_attempts}): {e}")
            
            if attempt < max_attempts - 1:
                print(f"Nouvelle tentative dans {wait_time} secondes...")
                time.sleep(wait_time)
                
                # Génération d'un nom unique pour éviter les conflits
                filepath = generate_unique_filename(filepath)
                print(f"Tentative avec un nouveau nom de fichier: {filepath}")
            else:
                print("Échec après plusieurs tentatives.")
                raise
        except Exception as e:
            print(f"Erreur inattendue lors de la sauvegarde: {e}")
            raise

def add_hyperlink(slide, text, url, left, top, width, height):
    """Ajoute un hyperlien à une diapositive."""
    # Création d'une forme de texte
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    
    # Ajout du texte avec lien
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    
    # Formatage du texte
    run.font.color.rgb = ACCENT_COLOR
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.underline = True
    
    # Ajout du lien hypertexte
    p.hyperlink.address = url
    
    return textbox

def create_presentation():
    """Crée une présentation concise de max 5 slides."""
    print("Création de la présentation PowerPoint globale (5 slides maximum)...")
    print(f"Utilisation du modèle: {TEMPLATE_FILE}")
    
    # Vérifier si le modèle existe
    if not os.path.exists(TEMPLATE_FILE):
        print(f"ERREUR: Le fichier modèle {TEMPLATE_FILE} n'existe pas!")
        print("Veuillez vérifier le chemin du modèle et réessayer.")
        return None
    
    # Création de la présentation à partir du modèle
    prs = Presentation(TEMPLATE_FILE)
    
    # ----- SLIDE 1: Page de titre -----
    slide_layout = prs.slide_layouts[0]  # Layout de titre
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre principal
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "n8n Docker HTTPS Setup"
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.color.rgb = TITLE_COLOR
            paragraph.font.size = Pt(44)
            paragraph.font.bold = True
    
    # Sous-titre
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:  # Généralement le sous-titre
            shape.text = "Automatisation de workflows avec Docker, ngrok et IA"
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.color.rgb = TEXT_COLOR
                paragraph.font.size = Pt(28)
    
    # Ajouter le lien GitHub en bas de la diapositive
    add_hyperlink(slide, "Projet GitHub", GITHUB_LINK, 
                  Inches(0.5), Inches(6.5), Inches(11), Inches(0.5))
    
    # ----- SLIDE 2: Vue d'ensemble -----
    slide_layout = prs.slide_layouts[1]  # Layout avec titre et contenu
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Vue d'ensemble du projet"
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.color.rgb = TITLE_COLOR
            paragraph.font.size = Pt(40)
    
    # Contenu
    content_points = [
        "Plateforme n8n containerisée avec Docker",
        "Accès HTTPS sécurisé via tunneling ngrok",
        "Fonctionnalités IA et MCP intégrées",
        "Documentation interactive Streamlit",
        "Scripts utilitaires pour simplifier l'utilisation"
    ]
    
    # Trouver le placeholder pour le contenu
    content_shape = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:  # Corps du texte
            content_shape = shape
            break
    
    if content_shape:
        tf = content_shape.text_frame
        tf.clear()
        
        for i, point in enumerate(content_points):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = point
            p.font.color.rgb = TEXT_COLOR
            p.font.size = Pt(24)
            p.level = 0
    
    # ----- SLIDE 3: Architecture -----
    slide_layout = prs.slide_layouts[1]  # Layout avec titre et contenu
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Architecture technique"
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.color.rgb = TITLE_COLOR
            paragraph.font.size = Pt(40)
    
    # Contenu
    content_shape = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            content_shape = shape
            break
    
    if content_shape:
        tf = content_shape.text_frame
        tf.clear()
        
        architecture_points = [
            "Frontend: Interface n8n accessible via navigateur",
            "Backend: Container Docker avec l'image n8n officielle",
            "Sécurité: Tunnel HTTPS ngrok + Authentification basique",
            "Données: Volumes Docker persistants",
            "IA: Intégration avec Ollama ou OpenAI via MCP"
        ]
        
        for i, point in enumerate(architecture_points):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = point
            p.font.color.rgb = TEXT_COLOR
            p.font.size = Pt(24)
            p.level = 0
    
    # ----- SLIDE 4: Fonctionnalités clés -----
    slide_layout = prs.slide_layouts[2]  # Layout pour sections
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Fonctionnalités clés"
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.color.rgb = TITLE_COLOR
            paragraph.font.size = Pt(44)
            paragraph.font.bold = True
    
    # Ajouter des formes avec les fonctionnalités clés
    features = [
        ("Docker", "Environnement isolé et portable"),
        ("HTTPS", "Accès sécurisé via ngrok"),
        ("IA", "Workflows intelligents avec agents IA"),
        ("Documentation", "Interface Streamlit interactive")
    ]
    
    # Créer une grille 2x2 de fonctionnalités
    for i, (title, desc) in enumerate(features):
        left = Inches(1.5) if i % 2 == 0 else Inches(7)
        top = Inches(2) if i < 2 else Inches(4)
        
        # Ajouter une forme arrondie
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                     left, top, Inches(4), Inches(1.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_COLOR
        shape.line.color.rgb = RGBColor(255, 255, 255)
        
        # Ajouter texte
        tf = shape.text_frame
        tf.word_wrap = True
        
        # Titre de la fonctionnalité
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER
        
        # Description
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(10, 10, 10)
        p.alignment = PP_ALIGN.CENTER
    
    # ----- SLIDE 5: Conclusion et liens -----
    slide_layout = prs.slide_layouts[5]  # Layout final
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Accéder au projet"
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.color.rgb = TITLE_COLOR
            paragraph.font.size = Pt(44)
    
    # Ajouter des ressources et liens
    links = [
        ("Code source GitHub", GITHUB_LINK),
        ("Documentation Streamlit", "https://streamlit.io/cloud"),
        ("Site officiel n8n", "https://n8n.io/"),
        ("Docker Hub", "https://hub.docker.com/r/n8nio/n8n/")
    ]
    
    left = Inches(2.5)
    top = Inches(2.5)
    height = Inches(0.5)
    width = Inches(7)
    
    for title, url in links:
        textbox = add_hyperlink(slide, f"{title}: {url}", url, left, top, width, height)
        top += Inches(0.7)
    
    # Pied de page avec date
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(11), Inches(0.5))
    footer_p = footer.text_frame.add_paragraph()
    footer_p.text = f"Créé le {datetime.now().strftime('%d/%m/%Y')} | n8n Docker HTTPS Setup"
    footer_p.font.size = Pt(12)
    footer_p.font.color.rgb = TEXT_COLOR
    footer_p.alignment = PP_ALIGN.CENTER
    
    return prs

def main():
    """Fonction principale."""
    # Créer la présentation
    prs = create_presentation()
    if prs:
        # Sauvegarder avec gestion des erreurs
        saved_file = save_presentation_safely(prs, OUTPUT_FILE)
        print(f"Présentation globale créée avec succès: {saved_file}")
    else:
        print("Erreur lors de la création de la présentation.")

if __name__ == "__main__":
    main()
