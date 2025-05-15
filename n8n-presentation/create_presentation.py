import os
from datetime import datetime
import time
import random
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Configuration des chemins
OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_FILE = os.path.join(OUTPUT_DIR, "n8n_Docker_HTTPS_Setup.pptx")
TEMPLATE_FILE = os.path.join(r"d:\00-Conception_AI\__Automations_n8n\n8n-presentation", "ThinkinG_modele.pptx")

# Définition des couleurs pour le texte sur fond sombre
TITLE_COLOR = RGBColor(240, 240, 240)  # Blanc cassé pour les titres
TEXT_COLOR = RGBColor(220, 220, 220)   # Gris clair pour le texte principal
ACCENT_COLOR = RGBColor(86, 180, 233)  # Bleu clair pour l'accentuation

def create_title_slide(prs, title, subtitle):
    """Crée une diapositive de titre adaptée au modèle sombre."""
    slide_layout = prs.slide_layouts[0]  # Layout de titre
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre - adaptez ces lignes selon la structure réelle du modèle
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = title
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(44)
        title_paragraph.font.color.rgb = TITLE_COLOR
        title_paragraph.font.bold = True
    
    # Sous-titre - ajustez selon le modèle
    # Le placeholder peut être différent selon le modèle
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:  # Généralement le sous-titre
            shape.text = subtitle
            subtitle_paragraph = shape.text_frame.paragraphs[0]
            subtitle_paragraph.font.size = Pt(24)
            subtitle_paragraph.font.color.rgb = TEXT_COLOR
    
    return slide

def create_section_slide(prs, title):
    """Crée une diapositive de section adaptée au modèle sombre."""
    # Utilisez le layout dédié aux sections dans votre modèle
    # Le numéro peut varier, à ajuster en fonction de votre modèle
    slide_layout = prs.slide_layouts[2]  
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = title
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(40)
        title_paragraph.font.color.rgb = TITLE_COLOR
        title_paragraph.font.bold = True
    
    return slide

def create_content_slide(prs, title, content, bullet_points=True):
    """Crée une diapositive de contenu adaptée au modèle sombre."""
    slide_layout = prs.slide_layouts[1]  # Layout avec titre et contenu
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = title
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(36)
        title_paragraph.font.color.rgb = TITLE_COLOR
    
    # Trouver le placeholder pour le contenu
    content_shape = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:  # Généralement le corps du texte
            content_shape = shape
            break
    
    if content_shape:
        content_frame = content_shape.text_frame
        
        # Effacer le contenu par défaut
        if content_frame.paragraphs:
            p = content_frame.paragraphs[0]
            p.clear()
        else:
            p = content_frame.add_paragraph()
        
        if isinstance(content, list):
            # Premier point
            if content:
                p.text = content[0]
                p.font.color.rgb = TEXT_COLOR
                p.font.size = Pt(24)
                p.level = 0
            
            # Autres points
            for point in content[1:]:
                p = content_frame.add_paragraph()
                p.text = point
                p.font.color.rgb = TEXT_COLOR
                p.font.size = Pt(22)
                p.level = 0 if not bullet_points else 1
        else:
            p.text = content
            p.font.color.rgb = TEXT_COLOR
            p.font.size = Pt(24)
    
    return slide

def generate_unique_filename(base_path):
    """Génère un nom de fichier unique en cas de conflit."""
    if not os.path.exists(base_path):
        return base_path
    
    # Séparer le nom et l'extension
    directory, filename = os.path.split(base_path)
    name, ext = os.path.splitext(filename)
    
    # Créer un nouveau nom avec timestamp et nombre aléatoire
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    random_num = random.randint(1000, 9999)
    new_filename = f"{name}_{timestamp}_{random_num}{ext}"
    
    return os.path.join(directory, new_filename)

def save_presentation_safely(prs, filepath):
    """Sauvegarde la présentation avec gestion des erreurs de permission."""
    max_attempts = 3
    wait_time = 2  # secondes
    
    for attempt in range(max_attempts):
        try:
            if os.path.exists(filepath):
                print(f"Le fichier {filepath} existe déjà.")
                try:
                    # Tenter de le supprimer
                    os.remove(filepath)
                    print("Ancien fichier supprimé avec succès.")
                except PermissionError:
                    # Si on ne peut pas le supprimer, générer un nom unique
                    filepath = generate_unique_filename(filepath)
                    print(f"Impossible de supprimer l'ancien fichier. Utilisation d'un nouveau nom: {filepath}")
            
            # Essayer de sauvegarder la présentation
            prs.save(filepath)
            print(f"Présentation sauvegardée avec succès: {filepath}")
            return filepath
            
        except PermissionError as e:
            print(f"Erreur de permission (tentative {attempt+1}/{max_attempts}): {e}")
            
            if attempt < max_attempts - 1:
                print(f"Nouvelle tentative dans {wait_time} secondes...")
                time.sleep(wait_time)
                
                # Génération d'un nom unique pour éviter les conflits
                filepath = generate_unique_filename(filepath)
                print(f"Tentative avec un nouveau nom de fichier: {filepath}")
            else:
                print("Échec après plusieurs tentatives.")
                raise
        except Exception as e:
            print(f"Erreur inattendue lors de la sauvegarde: {e}")
            raise

def main():
    """Fonction principale pour créer la présentation."""
    print("Création de la présentation PowerPoint pour n8n Docker HTTPS Setup...")
    print(f"Utilisation du modèle: {TEMPLATE_FILE}")
    
    # Vérifier si le modèle existe
    if not os.path.exists(TEMPLATE_FILE):
        print(f"ERREUR: Le fichier modèle {TEMPLATE_FILE} n'existe pas!")
        print("Veuillez vérifier le chemin du modèle et réessayer.")
        return
    
    # Création de la présentation à partir du modèle
    prs = Presentation(TEMPLATE_FILE)
    
    # Le reste du code reste similaire, mais les fonctions d'ajout de diapositives
    # ont été adaptées pour le modèle sombre
    
    # Diapositive de titre
    create_title_slide(
        prs,
        "n8n Docker HTTPS Setup",
        f"Automatisation de workflows avec n8n, Docker et ngrok\n{datetime.now().strftime('%d/%m/%Y')}"
    )
    
    # Sections et diapositives de contenu comme avant
    # Continuer avec les mêmes sections mais avec le style adapté
    
    # Section 1: Introduction
    create_section_slide(prs, "Introduction")
    
    create_content_slide(
        prs,
        "Qu'est-ce que n8n?",
        [
            "n8n est une plateforme d'automatisation de workflows extensible",
            "Alternative open-source à Zapier ou Integromat",
            "Interface visuelle de création de workflows",
            "Auto-hébergeable et hautement personnalisable",
            "Supporte +200 intégrations avec des services externes"
        ]
    )
    
    create_content_slide(
        prs,
        "Objectifs du projet",
        [
            "Configurer n8n dans un environnement Docker",
            "Fournir un accès HTTPS sécurisé via ngrok",
            "Activer les fonctionnalités expérimentales d'IA",
            "Créer une documentation interactive avec Streamlit",
            "Partager facilement la configuration et les workflows"
        ]
    )
    
    # Section 2: Architecture
    create_section_slide(prs, "Architecture")
    
    create_content_slide(
        prs,
        "Composants du système",
        [
            "Docker: Conteneurisation de n8n",
            "Docker Compose: Orchestration des conteneurs",
            "ngrok: Tunneling HTTPS sécurisé",
            "n8n: Plateforme d'automatisation",
            "Ollama/MCP: Intégration IA locale",
            "Streamlit: Documentation interactive"
        ]
    )
    
    create_content_slide(
        prs,
        "Flux de données",
        [
            "1. L'utilisateur accède à n8n via l'URL ngrok HTTPS",
            "2. ngrok redirige le trafic vers le conteneur Docker",
            "3. n8n traite les workflows et s'intègre avec les services externes",
            "4. Les données persistantes sont stockées dans un volume Docker",
            "5. Les fonctionnalités d'IA utilisent Ollama ou des API externes"
        ]
    )
    
    # Section 3: Installation et Configuration
    create_section_slide(prs, "Installation et Configuration")
    
    create_content_slide(
        prs,
        "Prérequis",
        [
            "Docker et Docker Compose",
            "ngrok (compte gratuit)",
            "Python 3.8+ pour les scripts utilitaires",
            "Connexion Internet pour les images Docker",
            "Facultatif: Ollama pour l'IA locale"
        ]
    )
    
    create_content_slide(
        prs,
        "Configuration Docker",
        [
            "docker-compose.yml définit la configuration n8n",
            "Volumes persistants pour les données",
            "Variables d'environnement pour l'authentification",
            "Configuration des ports et de l'accès réseau",
            "Activation des fonctionnalités expérimentales d'IA"
        ]
    )
    
    create_content_slide(
        prs,
        "Configuration ngrok",
        [
            "Création d'un compte et obtention du token",
            "Configuration du token avec ngrok authtoken",
            "Tunnel vers le port n8n (5678)",
            "URL HTTPS publique générée automatiquement",
            "Interface web de surveillance sur http://localhost:4040"
        ]
    )
    
    # Section 4: Fonctionnalités IA
    create_section_slide(prs, "Fonctionnalités IA")
    
    create_content_slide(
        prs,
        "Activation de l'IA dans n8n",
        [
            "Variables d'environnement spécifiques:",
            "N8N_EXPERIMENTAL_MCP=true",
            "N8N_MCP_BACKEND_URL=http://localhost:11434",
            "Intégration avec Ollama ou autres backends",
            "Accès aux nodes LangChain et AI Agents"
        ]
    )
    
    create_content_slide(
        prs,
        "Workflow d'exemple: Agent IA pour le Scraping",
        [
            "Utilise un agent IA pour comprendre les demandes",
            "Exécute automatiquement des outils de scraping",
            "Combine chat, mémoire et exécution d'outils",
            "Interface utilisateur conversationnelle",
            "Extensible pour d'autres tâches d'IA"
        ]
    )
    
    # Section 5: Documentation
    create_section_slide(prs, "Documentation")
    
    create_content_slide(
        prs,
        "Documentation Streamlit",
        [
            "Interface web interactive pour la documentation",
            "Navigation claire entre les différentes sections",
            "Code source et exemples intégrés",
            "Déploiement sur le réseau local ou via ngrok",
            "Option de déploiement sur Streamlit Cloud"
        ]
    )
    
    create_content_slide(
        prs,
        "Export vers Notion",
        [
            "Extraction du contenu Markdown depuis Streamlit",
            "Conversion en format compatible avec Notion",
            "Guide détaillé pour l'importation",
            "Captures d'écran automatiques",
            "Structure organisée pour Notion"
        ]
    )
    
    # Section 6: Scripts utilitaires
    create_section_slide(prs, "Scripts utilitaires")
    
    create_content_slide(
        prs,
        "Scripts disponibles",
        [
            "check_prerequisites.py: Vérification des prérequis",
            "0-start_n8n_menu.py: Menu interactif de démarrage",
            "run_streamlit.bat: Lancement de la documentation",
            "export_to_notion.py: Export vers Notion",
            "create_presentation.py: Génération de cette présentation"
        ]
    )
    
    # Section 7: Conclusion
    create_section_slide(prs, "Conclusion")
    
    create_content_slide(
        prs,
        "Avantages du projet",
        [
            "Solution complète d'automatisation auto-hébergée",
            "Accès sécurisé depuis n'importe où via HTTPS",
            "Fonctionnalités d'IA intégrées",
            "Documentation riche et accessible",
            "Scripts utilitaires pour faciliter l'utilisation"
        ]
    )
    
    create_content_slide(
        prs,
        "Prochaines étapes",
        [
            "Ajout de nouveaux workflows d'exemple",
            "Intégration avec d'autres services AI/ML",
            "Amélioration de la sécurité",
            "Automatisation du déploiement",
            "Création d'une communauté d'utilisateurs"
        ]
    )
    
    # Diapositive finale
    slide_layout = prs.slide_layouts[5]  # Layout simple, à ajuster selon le modèle
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Merci!"
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(60)
        title_paragraph.font.color.rgb = TITLE_COLOR
        title_paragraph.alignment = PP_ALIGN.CENTER
    
    # Ajout d'un sous-titre
    left = Inches(1)
    top = Inches(4)
    width = Inches(8)
    height = Inches(1.5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    
    p = tf.add_paragraph()
    p.text = "n8n Docker HTTPS Setup"
    p.font.size = Pt(32)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "Documentation et code source disponibles sur GitHub"
    p.font.size = Pt(24)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Sauvegarde de la présentation avec gestion des erreurs
    saved_file = save_presentation_safely(prs, OUTPUT_FILE)
    
    print(f"Présentation créée avec succès: {saved_file}")
    print("Le modèle sombre a été utilisé avec des polices claires pour une meilleure lisibilité.")

if __name__ == "__main__":
    main()
